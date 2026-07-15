"""Generate the MediStruct project presentation (16 slides)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- Theme (matches the app) --------------------------------------------
TEAL      = RGBColor(0x0D, 0x94, 0x88)   # primary
TEAL_DK   = RGBColor(0x0B, 0x6B, 0x63)
INK       = RGBColor(0x0F, 0x17, 0x2A)   # foreground
MUTED     = RGBColor(0x64, 0x74, 0x8B)
LIGHT     = RGBColor(0xF6, 0xF8, 0xFB)   # background
CARD      = RGBColor(0xFF, 0xFF, 0xFF)
BORDER    = RGBColor(0xE2, 0xE8, 0xF0)
ACCENT    = RGBColor(0x4F, 0x46, 0xE5)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = LIGHT
    return s


def rect(slide, x, y, w, h, color, line=None, radius=False):
    from pptx.enum.shapes import MSO_SHAPE
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    sp = slide.shapes.add_shape(shape_type, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def txt(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=6, line_spacing=1.05):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    first = True
    for item in runs:
        text, size, color, bold = item[0], item[1], item[2], item[3]
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.name = "Segoe UI"
    return tb


def bullets(slide, x, y, w, h, items, size=16, color=INK, gap=10):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for it in items:
        lead, rest = (it if isinstance(it, tuple) else (it, ""))
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap)
        p.line_spacing = 1.08
        r = p.add_run()
        r.text = "▸ "
        r.font.size = Pt(size)
        r.font.color.rgb = TEAL
        r.font.bold = True
        r.font.name = "Segoe UI"
        rb = p.add_run()
        rb.text = lead
        rb.font.size = Pt(size)
        rb.font.color.rgb = INK
        rb.font.bold = True
        rb.font.name = "Segoe UI"
        if rest:
            rr = p.add_run()
            rr.text = " — " + rest
            rr.font.size = Pt(size)
            rr.font.color.rgb = MUTED
            rr.font.bold = False
            rr.font.name = "Segoe UI"
    return tb


def header(slide, kicker, title, num):
    rect(slide, 0, 0, Inches(0.28), SH, TEAL)
    txt(slide, Inches(0.7), Inches(0.45), Inches(10), Inches(0.4),
        [(kicker.upper(), 13, TEAL, True)])
    txt(slide, Inches(0.7), Inches(0.78), Inches(11.5), Inches(0.9),
        [(title, 30, INK, True)])
    rect(slide, Inches(0.72), Inches(1.62), Inches(1.1), Pt(3), TEAL)
    # slide number chip
    txt(slide, Inches(12.2), Inches(6.95), Inches(1), Inches(0.4),
        [(f"{num:02d} / 16", 11, MUTED, False)], align=PP_ALIGN.RIGHT)


def card(slide, x, y, w, h, title, body_items, tint=CARD, accent=TEAL):
    c = rect(slide, x, y, w, h, tint, line=BORDER, radius=True)
    rect(slide, x, y, Inches(0.09), h, accent, radius=False)
    txt(slide, x + Inches(0.3), y + Inches(0.22), w - Inches(0.5), Inches(0.5),
        [(title, 16, INK, True)])
    tb = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.72),
                                  w - Inches(0.55), h - Inches(0.9))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for it in body_items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(5)
        p.line_spacing = 1.05
        r = p.add_run()
        r.text = it
        r.font.size = Pt(12.5)
        r.font.color.rgb = MUTED
        r.font.name = "Segoe UI"
    return c


def table(slide, rows):
    """3-column styled table: Structure | Real feature | Complexity."""
    x = Inches(0.7)
    y = Inches(2.05)
    w = Inches(11.95)
    col1 = Inches(3.1)
    col3 = Inches(3.2)
    col2 = Emu(int(w) - int(col1) - int(col3))
    rh = Inches(0.55)
    # header row
    rect(slide, x, y, w, rh, TEAL)
    hx = [x, Emu(int(x) + int(col1)), Emu(int(x) + int(col1) + int(col2))]
    for cx, cw, label in zip(hx, [col1, col2, col3],
                             ["Structure", "Real feature", "Complexity"]):
        txt(slide, Emu(int(cx) + int(Inches(0.2))), y, cw, rh,
            [(label, 13, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)
    # body rows
    for i, (name, use, comp) in enumerate(rows):
        ry = Emu(int(y) + int(rh) + i * int(Inches(0.68)))
        band = CARD if i % 2 == 0 else RGBColor(0xEE, 0xF3, 0xF7)
        rect(slide, x, ry, w, Inches(0.68), band, line=BORDER)
        txt(slide, Emu(int(hx[0]) + int(Inches(0.2))), ry, col1, Inches(0.68),
            [(name, 13.5, INK, True)], anchor=MSO_ANCHOR.MIDDLE)
        txt(slide, Emu(int(hx[1]) + int(Inches(0.2))), ry, col2, Inches(0.68),
            [(use, 12.5, MUTED, False)], anchor=MSO_ANCHOR.MIDDLE)
        txt(slide, Emu(int(hx[2]) + int(Inches(0.2))), ry, col3, Inches(0.68),
            [(comp, 12, TEAL_DK, True)], anchor=MSO_ANCHOR.MIDDLE)


# =========================================================================
# Slide 1 — Title / Introduction
# =========================================================================
s = add_slide()
rect(s, 0, 0, SW, SH, INK)
rect(s, 0, 0, SW, Inches(2.5), TEAL_DK)
rect(s, 0, Inches(2.5), SW, Inches(0.06), TEAL)
# logo mark
lg = rect(s, Inches(0.9), Inches(0.85), Inches(0.8), Inches(0.8), TEAL, radius=True)
txt(s, Inches(0.9), Inches(0.92), Inches(0.8), Inches(0.7), [("✚", 34, WHITE, True)],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
txt(s, Inches(1.9), Inches(0.95), Inches(9), Inches(0.7), [("MediStruct", 30, WHITE, True)],
    anchor=MSO_ANCHOR.MIDDLE)
txt(s, Inches(1.9), Inches(1.55), Inches(10), Inches(0.5),
    [("Hospital Management System · DSA Project", 15, RGBColor(0xCC,0xFB,0xF1), False)])

txt(s, Inches(0.9), Inches(3.0), Inches(11.5), Inches(1.2),
    [("A Hospital Management System where every", 30, WHITE, True),
     ("workflow is powered by a hand-written data structure.", 30, WHITE, True)],
    line_spacing=1.1)
txt(s, Inches(0.9), Inches(4.35), Inches(11.5), Inches(0.6),
    [("Registration → specialist routing → triage → consultation → labs → operation → admission → discharge",
      14, MUTED, False)])

# Team card
rect(s, Inches(0.9), Inches(5.15), Inches(11.5), Inches(1.55), RGBColor(0x1E,0x29,0x3B), radius=True)
txt(s, Inches(1.2), Inches(5.32), Inches(4), Inches(0.4), [("PRESENTED BY", 12, TEAL, True)])
names = ["Rijja", "Shamim", "Ayesha", "Muzzamil"]
nx = Inches(1.2)
gap = Inches(2.75)
for i, n in enumerate(names):
    x = Emu(int(nx) + i * int(gap))
    dot = rect(s, x, Inches(5.75), Inches(0.42), Inches(0.42), TEAL, radius=True)
    txt(s, x, Inches(5.79), Inches(0.42), Inches(0.36), [(n[0], 16, WHITE, True)],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Emu(int(x) + int(Inches(0.55))), Inches(5.78), Inches(2.1), Inches(0.4),
        [(n, 17, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)

# =========================================================================
# Slide 2 — Architecture (two-layer)
# =========================================================================
s = add_slide()
header(s, "Architecture", "Two-Layer Design", 2)
txt(s, Inches(0.7), Inches(1.85), Inches(12), Inches(0.6),
    [("The structures do real work instead of being decoration:", 16, MUTED, False)])
card(s, Inches(0.7), Inches(2.7), Inches(5.9), Inches(2.0), "PostgreSQL (Neon)",
     ["Permanent storage — the single source of truth.",
      "Patients, doctors, bills, appointments, wards, events.",
      "Accessed through server actions & typed queries."])
card(s, Inches(6.75), Inches(2.7), Inches(5.9), Inches(2.0),
     "Hand-written DSA Layer  (src/lib/dsa)",
     ["The in-memory working layer.",
      "Rows are loaded, fed through a Queue / Graph / Heap / Sort,",
      "then the result is rendered."], accent=ACCENT)
# flow strip
rect(s, Inches(0.7), Inches(5.1), Inches(11.95), Inches(1.15), CARD, line=BORDER, radius=True)
flow = "Read rows from Postgres   →   Build a data structure   →   Run its algorithm   →   Render the result"
txt(s, Inches(0.7), Inches(5.1), Inches(11.95), Inches(1.15), [(flow, 17, TEAL_DK, True)],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# =========================================================================
# Slide 3 — Tech Stack
# =========================================================================
s = add_slide()
header(s, "Foundation", "Technology Stack", 3)
stack = [
    ("Next.js 16", "App Router, React Server Components & server actions"),
    ("TypeScript", "End-to-end type safety across UI, actions and data"),
    ("Tailwind CSS v4", "Utility-first styling; responsive & theme-aware"),
    ("Better Auth", "Email/password auth with role-based access control"),
    ("Neon Postgres", "Serverless Postgres — the permanent data store"),
    ("Hand-written DSA", "src/lib/dsa — every structure built from scratch"),
]
gx, gy = Inches(0.7), Inches(2.0)
cw, ch = Inches(3.85), Inches(1.75)
for i, (t, d) in enumerate(stack):
    r, c = divmod(i, 3)
    x = Emu(int(gx) + c * int(Inches(4.04)))
    y = Emu(int(gy) + r * int(Inches(1.95)))
    card(s, x, y, cw, ch, t, [d])

# =========================================================================
# Slide 4 — User Roles
# =========================================================================
s = add_slide()
header(s, "Access Control", "Three User Roles", 4)
roles = [
    ("Administrator", TEAL,
     ["Manages doctors & receptionists.",
      "Configures wards & beds.",
      "Sees system-wide dashboard."]),
    ("Receptionist", ACCENT,
     ["Registers patients & sets triage severity.",
      "Routes issue → specialist, assigns doctor.",
      "Handles billing, insurance & the waiting queue."]),
    ("Doctor", RGBColor(0xD9,0x77,0x06),
     ["Works a priority queue (most-critical first).",
      "Writes clinical notes; orders labs & operations.",
      "Recommends admission / discharge."]),
]
cw, ch = Inches(3.85), Inches(4.0)
for i, (name, col, items) in enumerate(roles):
    x = Emu(int(Inches(0.7)) + i * int(Inches(4.04)))
    rect(s, x, Inches(2.1), cw, ch, CARD, line=BORDER, radius=True)
    rect(s, x, Inches(2.1), cw, Inches(0.85), col, radius=True)
    txt(s, x, Inches(2.22), cw, Inches(0.6), [(name, 19, WHITE, True)],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    bullets(s, Emu(int(x)+int(Inches(0.3))), Inches(3.2), Emu(int(cw)-int(Inches(0.55))),
            Inches(2.7), items, size=14)

# =========================================================================
# Slide 5 — Patient Journey
# =========================================================================
s = add_slide()
header(s, "The Big Picture", "End-to-End Patient Journey", 5)
steps = ["Register", "Route to\nSpecialist", "Triage\nSeverity", "Pay\nConsult",
         "Doctor\nQueue", "Consult\n& Notes", "Labs /\nOperation", "Admit /\nDischarge"]
n = len(steps)
bw = Inches(1.42)
gapx = Inches(0.08)
total = int(bw) * n + int(gapx) * (n - 1)
startx = int((int(SW) - total) / 2)
y = Inches(3.1)
for i, st in enumerate(steps):
    x = Emu(startx + i * (int(bw) + int(gapx)))
    col = TEAL if i % 2 == 0 else TEAL_DK
    rect(s, x, y, bw, Inches(1.3), col, radius=True)
    txt(s, x, y, bw, Inches(1.3), [(st, 12.5, WHITE, True)],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    if i < n - 1:
        ax = Emu(x + int(bw) - int(Inches(0.02)))
        txt(s, Emu(int(ax)), Emu(int(y)+int(Inches(0.45))), Inches(0.12), Inches(0.4),
            [("›", 20, MUTED, True)], align=PP_ALIGN.CENTER)
txt(s, Inches(0.7), Inches(4.9), Inches(12), Inches(1.2),
    [("Every action appends a row to patient_events — an immutable audit trail that becomes the "
      "patient's timeline, walkable forward & back (Doubly Linked List).", 15, MUTED, False)],
    align=PP_ALIGN.CENTER, line_spacing=1.15)

# =========================================================================
# Module slides helper
# =========================================================================
def module_slide(num, kicker, title, intro, left_title, left, right_title, right, ds):
    s = add_slide()
    header(s, kicker, title, num)
    txt(s, Inches(0.7), Inches(1.8), Inches(12), Inches(0.65),
        [(intro, 15.5, MUTED, False)], line_spacing=1.12)
    card(s, Inches(0.7), Inches(2.75), Inches(5.9), Inches(3.0), left_title, left)
    card(s, Inches(6.75), Inches(2.75), Inches(5.9), Inches(3.0), right_title, right, accent=ACCENT)
    rect(s, Inches(0.7), Inches(5.95), Inches(11.95), Inches(0.75), RGBColor(0xCC,0xFB,0xF1), radius=True)
    txt(s, Inches(1.0), Inches(5.95), Inches(11.4), Inches(0.75),
        [("DSA used:  " + ds, 14, TEAL_DK, True)], anchor=MSO_ANCHOR.MIDDLE)
    return s

# Slide 6 — Reception module
module_slide(6, "Module 01", "Reception & Registration",
    "The front door of the hospital: every patient starts here.",
    "What it does",
    ["Registers a patient (auto integer ID from a Postgres serial).",
     "Captures issue, contact & payment details.",
     "Sets a triage severity (1 critical … 5 routine).",
     "Opens the consultation bill automatically."],
    "Waiting queue",
    ["Reception's waiting line is first-come, first-served.",
     "Loaded from Postgres by arrival time.",
     "Fed through the hand-written FIFO Queue.",
     "Front of the line is shown as 'serve next'."],
    "Queue (FIFO)  ·  Dynamic Array  ·  Hash Map")

# Slide 7 — Specialist Routing
module_slide(7, "Module 02", "Specialist Routing (Graph)",
    "The most important algorithmic feature — mapping a symptom to the right doctor.",
    "How routing works",
    ["Symptoms & specialties are nodes in a directed weighted graph.",
     "Edges (symptom → specialty) carry a routing weight.",
     "Dijkstra's shortest path finds the best department.",
     "Rebuilt each request from the symptom_routes table."],
    "Ranking the doctors",
    ["The department's doctors are then ranked.",
     "Our own quickSort orders them by rating, then fee.",
     "Receptionist sees the best-matching specialist first.",
     "Manual override is always available."],
    "Graph (adjacency list) + Dijkstra / BFS  ·  quickSort")

# Slide 8 — Emergency Triage (highlight)
s = add_slide()
header(s, "Module 03 · Flagship Feature", "Emergency Triage (Priority Queue)", 8)
txt(s, Inches(0.7), Inches(1.8), Inches(12), Inches(0.65),
    [("A critical patient must be able to jump ahead of a routine one — a plain FIFO queue cannot do this.",
      15.5, MUTED, False)], line_spacing=1.12)
card(s, Inches(0.7), Inches(2.7), Inches(5.9), Inches(2.55), "How it works",
     ["Reception sets each patient's severity at registration.",
      "The doctor's queue is a Min Priority Queue (binary heap).",
      "Most-critical (lowest severity number) is dequeued first — O(log n).",
      "Order is independent of arrival time."])
card(s, Inches(6.75), Inches(2.7), Inches(5.9), Inches(2.55), "Fair tie-breaking",
     ["Same severity ⇒ longest-waiting patient goes first (FIFO fallback).",
      "Composite priority = severity × 1,000,000 + arrivalIndex.",
      "Severity always dominates; arrival only breaks ties."], accent=ACCENT)
rect(s, Inches(0.7), Inches(5.45), Inches(11.95), Inches(1.25), INK, radius=True)
txt(s, Inches(1.0), Inches(5.58), Inches(11.4), Inches(0.4),
    [("Example — triage order", 12, TEAL, True)])
txt(s, Inches(1.0), Inches(5.9), Inches(11.4), Inches(0.7),
    [("Ali (sev 3, arrived last)  →  Bilal (sev 4)  →  Sana (sev 4)  →  Muzzamil (sev 5, arrived first)",
      15, WHITE, True)], anchor=MSO_ANCHOR.MIDDLE)

# Slide 9 — Doctor module
module_slide(9, "Module 04", "Doctor Workspace",
    "Where the clinician sees patients and records care.",
    "My Queue",
    ["Shows patients assigned & waiting to be seen.",
     "Ordered by triage severity (priority queue).",
     "'See next' marks the true triage front."],
    "Clinical console",
    ["Writes chronological notes for a patient.",
     "Notes stored & replayed via a Singly Linked List.",
     "Orders lab tests & recommends operations / admission."],
    "Min Priority Queue  ·  Singly Linked List")

# Slide 10 — Billing & Insurance
module_slide(10, "Module 05", "Billing & Insurance",
    "Accurate charges for consultations, labs, operations and admissions.",
    "Billing rules",
    ["Every chargeable item becomes a bill row.",
     "Concession applies to lab/operation/admission — never consultation.",
     "Order: concession first, then insurance %, then patient pays rest.",
     "Bill line-items modelled with a Dynamic Array."],
    "Insurance & operations",
    ["Percentage coverage per patient (e.g. 80%).",
     "Operation loop: schedule → must pay before date → else reschedule.",
     "Bills indexed / searched by ID (binary search over sorted IDs)."],
    "Dynamic Array + binary search  ·  Sorting")

# Slide 11 — Admin module
module_slide(11, "Module 06", "Administration",
    "System configuration and staff management.",
    "People",
    ["Create & manage doctors (specialty, fee, rating).",
     "Create & manage receptionists.",
     "Role-based access enforced by Better Auth."],
    "Wards, beds & dashboard",
    ["Define wards and their beds.",
     "Admitting a patient occupies a bed; discharge frees it.",
     "Dashboard summarises patients, staff and occupancy."],
    "Hash Map (O(1) lookups)  ·  Red-Black Tree (date index)")

# Slide 12 — DSA Gallery
s = add_slide()
header(s, "Proof", "The DSA Gallery (/dsa)", 12)
txt(s, Inches(0.7), Inches(1.85), Inches(12), Inches(0.7),
    [("A dedicated page that proves the theory is real — interactive demos plus a live Big-O panel for every structure.",
      16, MUTED, False)], line_spacing=1.15)
card(s, Inches(0.7), Inches(2.85), Inches(3.85), Inches(3.4), "Triage demo",
     ["Runs the real Min Priority Queue in the browser.",
      "'Serve next' always picks the most-critical patient.",
      "Now backs the live doctor queue too."])
card(s, Inches(4.74), Inches(2.85), Inches(3.85), Inches(3.4), "Sort race",
     ["Runs all 5 hand-written sorts on your numbers.",
      "Counts real comparisons & swaps.",
      "Ranks algorithms fewest-comparisons first."], accent=ACCENT)
card(s, Inches(8.78), Inches(2.85), Inches(3.85), Inches(3.4), "Complexity cards",
     ["One card per structure.",
      "Big-O per operation + space cost.",
      "States the exact hospital feature it powers."])

# Slide 13 — Core data structures
s = add_slide()
header(s, "Under the Hood I", "Core Data Structures", 13)
rows = [
    ("Dynamic Array", "Bill line-items; binary search over IDs", "O(1) access · O(log n) search"),
    ("Queue (FIFO)", "Reception waiting line", "O(1) enqueue / dequeue"),
    ("Stack (LIFO)", "Undo / action history", "O(1) push / pop"),
    ("Min Priority Queue", "Emergency triage — most-critical first", "O(log n) enqueue / dequeue"),
    ("Singly Linked List", "Doctor's chronological note log", "O(1) append · O(n) find"),
    ("Doubly Linked List", "Patient timeline — step forward/back", "O(1) back / forward"),
]
table(s, rows)

# Slide 14 — Advanced structures
s = add_slide()
header(s, "Under the Hood II", "Trees, Hashing, Graphs & Sorting", 14)
rows2 = [
    ("Binary Search Tree", "Patient lookup by auto-generated ID", "O(log n) avg search"),
    ("Red-Black Tree", "Balanced appointment index by date", "O(log n) guaranteed"),
    ("Hash Map (chaining)", "O(1) patient / insurance / doctor lookup", "O(1) average get / set"),
    ("Graph + Dijkstra", "Symptom → specialist routing", "O(V²) shortest path"),
    ("Sorting × 5", "Rank doctors; live sort-race", "Merge/Quick: O(n log n)"),
]
table(s, rows2)

# Slide 15 — Conclusion
s = add_slide()
rect(s, 0, 0, SW, SH, INK)
rect(s, 0, 0, SW, Inches(0.06), TEAL)
txt(s, Inches(0.9), Inches(0.8), Inches(6), Inches(0.4), [("CONCLUSION", 13, TEAL, True)])
txt(s, Inches(0.9), Inches(1.15), Inches(11.5), Inches(1.0),
    [("A real product, powered end-to-end by DSA.", 32, WHITE, True)])
takeaways = [
    ("Every workflow maps to a data structure", "queues, heaps, trees, graphs, hashing & sorting — all hand-written."),
    ("Two-layer design keeps DSA doing real work", "Postgres stores truth; the DSA layer computes each feature."),
    ("Emergency triage is a live, fair priority queue", "severity-ordered, with FIFO tie-breaking."),
    ("The DSA Gallery proves it interactively", "live demos + Big-O for every structure."),
]
ty = Inches(2.4)
for i, (t, d) in enumerate(takeaways):
    y = Emu(int(ty) + i * int(Inches(0.85)))
    rect(s, Inches(0.9), y, Inches(0.12), Inches(0.62), TEAL, radius=True)
    txt(s, Inches(1.25), y, Inches(11), Inches(0.3), [(t, 17, WHITE, True)])
    txt(s, Inches(1.25), Emu(int(y)+int(Inches(0.33))), Inches(11), Inches(0.3),
        [(d, 13, MUTED, False)])
# thank you + names
rect(s, Inches(0.9), Inches(6.15), Inches(11.5), Inches(0.9), RGBColor(0x1E,0x29,0x3B), radius=True)
txt(s, Inches(1.2), Inches(6.15), Inches(3), Inches(0.9), [("Thank you", 20, WHITE, True)],
    anchor=MSO_ANCHOR.MIDDLE)
txt(s, Inches(4.2), Inches(6.15), Inches(8), Inches(0.9),
    [("Rijja  ·  Shamim  ·  Ayesha  ·  Muzzamil", 16, RGBColor(0xCC,0xFB,0xF1), True)],
    anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)

# =========================================================================
# Slide 16 — Speaker assignments / roadmap
# =========================================================================
s = add_slide()
header(s, "Who Presents What", "Presentation Roadmap", 16)
txt(s, Inches(0.7), Inches(1.72), Inches(12), Inches(0.45),
    [("How our four-person team divides the talk:", 15, MUTED, False)])

members = [
    ("Rijja", TEAL, "Stack & Roles", "Slides 3–4",
     ["Technology stack: Next.js, TypeScript, Tailwind, Better Auth, Neon.",
      "The three user roles & role-based access control."]),
    ("Shamim", ACCENT, "System Flow", "Slides 2, 5–7",
     ["Two-layer architecture (Postgres + hand-written DSA).",
      "End-to-end patient journey.",
      "Reception, registration & graph-based specialist routing."]),
    ("Ayesha", RGBColor(0xD9, 0x77, 0x06), "Core Modules & DSA", "Slides 8–11",
     ["Emergency triage — the live priority queue.",
      "Doctor workspace, billing & administration.",
      "The data structure behind each module."]),
    ("Muzzamil", TEAL_DK, "DSA Deep-Dive", "Slides 12–14",
     ["The interactive DSA Gallery.",
      "Core & advanced data structures.",
      "Complexity (Big-O) analysis of each."]),
]
positions = [(Inches(0.7), Inches(2.25)), (Inches(6.75), Inches(2.25)),
             (Inches(0.7), Inches(4.55)), (Inches(6.75), Inches(4.55))]
cw, ch = Inches(5.9), Inches(2.15)
for (name, col, topic, rng, items), (x, y) in zip(members, positions):
    rect(s, x, y, cw, ch, CARD, line=BORDER, radius=True)
    rect(s, x, y, Inches(0.09), ch, col)
    rect(s, x + Inches(0.3), y + Inches(0.28), Inches(0.62), Inches(0.62), col, radius=True)
    txt(s, x + Inches(0.3), y + Inches(0.32), Inches(0.62), Inches(0.54),
        [(name[0], 22, WHITE, True)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x + Inches(1.1), y + Inches(0.26), Inches(3.0), Inches(0.4), [(name, 19, INK, True)])
    txt(s, x + Inches(1.1), y + Inches(0.68), Inches(3.4), Inches(0.35), [(topic, 13, col, True)])
    txt(s, x + Inches(3.9), y + Inches(0.3), Inches(1.85), Inches(0.4), [(rng, 12, MUTED, True)],
        align=PP_ALIGN.RIGHT)
    bullets(s, x + Inches(0.32), y + Inches(1.14), cw - Inches(0.6), Inches(0.95),
            items, size=12, gap=4)

prs.save("MediStruct-Presentation.pptx")
print("saved MediStruct-Presentation.pptx with", len(prs.slides._sldIdLst), "slides")
